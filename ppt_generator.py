"""
PPT Generator - 使用 python-pptx 生成专业演示文稿
增强版：包含图标、卡片式布局、装饰元素
"""
import os
import uuid
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================
# 图标形状映射（用于要点前的图标）
# ============================================
POINT_ICONS = [
    MSO_SHAPE.DIAMOND,           # 菱形
    MSO_SHAPE.RIGHT_ARROW,       # 箭头
    MSO_SHAPE.STAR_5_POINT,      # 五角星
    MSO_SHAPE.LIGHTNING_BOLT,    # 闪电
    MSO_SHAPE.HEART,             # 心形
    MSO_SHAPE.OVAL,              # 圆形
    MSO_SHAPE.CHEVRON,           # V形
    MSO_SHAPE.CLOUD,             # 云形
]

# 页码图标
PAGE_ICONS = [
    MSO_SHAPE.PENTAGON,          # 五边形
]

# Color schemes for different templates
COLOR_SCHEMES = {
    "blueprint": {
        "primary": RGBColor(0x1A, 0x3C, 0x6E),
        "accent": RGBColor(0x3B, 0x82, 0xF6),
        "accent2": RGBColor(0x60, 0xA5, 0xFA),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "text": RGBColor(0x1E, 0x29, 0x3B),
        "text_light": RGBColor(0x64, 0x74, 0x8B),
        "light": RGBColor(0xEF, 0xF6, 0xFF),
        "card_bg": RGBColor(0xF8, 0xFA, 0xFC),
    },
    "structure": {
        "primary": RGBColor(0x06, 0x5F, 0x46),
        "accent": RGBColor(0x10, 0xB9, 0x81),
        "accent2": RGBColor(0x34, 0xD3, 0x99),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "text": RGBColor(0x1E, 0x29, 0x3B),
        "text_light": RGBColor(0x64, 0x74, 0x8B),
        "light": RGBColor(0xEC, 0xFD, 0xF5),
        "card_bg": RGBColor(0xF0, 0xFD, 0xF4),
    },
    "storytelling": {
        "primary": RGBColor(0x7C, 0x2D, 0x12),
        "accent": RGBColor(0xF5, 0x9E, 0x0B),
        "accent2": RGBColor(0xFB, 0xBF, 0x24),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "text": RGBColor(0x1E, 0x29, 0x3B),
        "text_light": RGBColor(0x64, 0x74, 0x8B),
        "light": RGBColor(0xFF, 0xFB, 0xEB),
        "card_bg": RGBColor(0xFF, 0xF7, 0xED),
    },
    "visual": {
        "primary": RGBColor(0x5B, 0x21, 0xB6),
        "accent": RGBColor(0x8B, 0x5C, 0xF6),
        "accent2": RGBColor(0xA7, 0x8B, 0xFA),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "text": RGBColor(0x1E, 0x29, 0x3B),
        "text_light": RGBColor(0x64, 0x74, 0x8B),
        "light": RGBColor(0xF5, 0xF3, 0xFF),
        "card_bg": RGBColor(0xFA, 0xF5, 0xFF),
    },
    "content": {
        "primary": RGBColor(0x0E, 0x74, 0x90),
        "accent": RGBColor(0x06, 0xB6, 0xD4),
        "accent2": RGBColor(0x22, 0xD3, 0xEE),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "text": RGBColor(0x1E, 0x29, 0x3B),
        "text_light": RGBColor(0x64, 0x74, 0x8B),
        "light": RGBColor(0xEC, 0xFE, 0xFF),
        "card_bg": RGBColor(0xF0, 0xFD, 0xFA),
    },
    "clarity": {
        "primary": RGBColor(0xBE, 0x12, 0x3C),
        "accent": RGBColor(0xF4, 0x3F, 0x5E),
        "accent2": RGBColor(0xFB, 0x71, 0x85),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "text": RGBColor(0x1E, 0x29, 0x3B),
        "text_light": RGBColor(0x64, 0x74, 0x8B),
        "light": RGBColor(0xFF, 0xF1, 0xF2),
        "card_bg": RGBColor(0xFF, 0xF5, 0xF5),
    },
}


def _add_rect(slide, color, left, top, width, height, radius=None):
    """添加矩形/圆角矩形"""
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        # 调整圆角大小
        shape.adjustments[0] = radius
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_circle(slide, color, left, top, size):
    """添加圆形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_icon_shape(slide, shape_type, color, left, top, size):
    """添加图标形状"""
    shape = slide.shapes.add_shape(shape_type, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _set_text(text_frame, text, font_size=18, color=None, bold=False, alignment=PP_ALIGN.LEFT):
    """设置文本框内容"""
    text_frame.clear()
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return p


def _set_multiline_text(text_frame, lines, font_size=16, color=None, bold=False, alignment=PP_ALIGN.LEFT, spacing=Pt(6)):
    """设置多行文本"""
    text_frame.clear()
    text_frame.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.alignment = alignment
        p.space_after = spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color


def generate_title_slide(prs, title, subtitle, colors):
    """生成增强版标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 主背景
    _add_rect(slide, colors["primary"], Inches(0), Inches(0), Inches(13.333), Inches(7.5))

    # 装饰性大圆（右上角）
    _add_circle(slide, colors["accent"], Inches(9.5), Inches(-0.5), Inches(3.5))

    # 装饰性小圆（右上角叠加）
    circle2 = _add_circle(slide, colors["accent2"], Inches(10.8), Inches(0.3), Inches(1.8))
    circle2.fill.fore_color.brightness = 0.2

    # 装饰性圆点（左侧）
    for i in range(5):
        dot = _add_circle(slide, colors["accent"], Inches(0.3 + i * 0.25), Inches(6.8), Inches(0.12))
        dot.fill.fore_color.brightness = 0.3 + i * 0.1

    # 顶部装饰条
    _add_rect(slide, colors["accent"], Inches(0), Inches(0), Inches(13.333), Inches(0.08))

    # 左侧装饰色块
    _add_rect(slide, colors["accent"], Inches(0.6), Inches(2.0), Inches(0.08), Inches(2.5))

    # 标题文字
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(8.5), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    _set_text(tf, title, font_size=42, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)

    # 副标题
    txBox2 = slide.shapes.add_textbox(Inches(1.0), Inches(4.3), Inches(8.5), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    _set_text(tf2, subtitle, font_size=16, color=RGBColor(0xCB, 0xD5, 0xE1))

    # 底部装饰条
    _add_rect(slide, colors["accent"], Inches(1.0), Inches(4.0), Inches(2.5), Inches(0.06))

    # 底部信息栏
    _add_rect(slide, colors["primary"], Inches(0), Inches(6.8), Inches(13.333), Inches(0.7))
    txInfo = slide.shapes.add_textbox(Inches(1.0), Inches(6.85), Inches(11), Inches(0.5))
    tfInfo = txInfo.text_frame
    _set_text(tfInfo, "Powered by MiMo AI  ·  智能演示文稿生成器",
              font_size=11, color=RGBColor(0x94, 0xA3, 0xB8), alignment=PP_ALIGN.CENTER)

    return slide


def generate_toc_slide(prs, slides, colors):
    """生成目录页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 左侧色块
    _add_rect(slide, colors["primary"], Inches(0), Inches(0), Inches(4.0), Inches(7.5))

    # 左侧标题
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(3.0), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    _set_text(tf, "目录", font_size=36, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)

    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(3.0), Inches(1))
    tf2 = txBox2.text_frame
    _set_text(tf2, "CONTENTS", font_size=14, color=RGBColor(0xCB, 0xD5, 0xE1))

    # 装饰圆
    _add_circle(slide, colors["accent"], Inches(2.8), Inches(0.5), Inches(1.2))

    # 右侧目录列表
    start_y = 0.8
    for i, s in enumerate(slides[:8]):  # 最多显示8项
        y = start_y + i * 0.75

        # 编号圆
        _add_circle(slide, colors["accent"], Inches(4.5), Inches(y + 0.05), Inches(0.45))
        txNum = slide.shapes.add_textbox(Inches(4.5), Inches(y + 0.05), Inches(0.45), Inches(0.45))
        tfNum = txNum.text_frame
        tfNum.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tfNum.paragraphs[0].add_run()
        run.text = str(i + 1)
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # 标题文字
        txTitle = slide.shapes.add_textbox(Inches(5.2), Inches(y), Inches(7.5), Inches(0.55))
        tfTitle = txTitle.text_frame
        tfTitle.word_wrap = True
        _set_text(tfTitle, s.get("title", f"第 {i+1} 页"), font_size=16, color=colors["text"])

        # 分隔线
        if i < len(slides[:8]) - 1:
            _add_rect(slide, colors["card_bg"], Inches(5.2), Inches(y + 0.65), Inches(7.0), Inches(0.01))

    return slide


def generate_content_slide_v2(prs, slide_data, slide_index, total_slides, colors):
    """生成增强版内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 左侧装饰条
    _add_rect(slide, colors["primary"], Inches(0), Inches(0), Inches(0.06), Inches(7.5))

    # 顶部浅色背景区域
    _add_rect(slide, colors["light"], Inches(0), Inches(0), Inches(13.333), Inches(1.5))

    # 标题图标色块
    _add_rect(slide, colors["primary"], Inches(0.6), Inches(0.35), Inches(0.5), Inches(0.5), radius=0.15)
    txIcon = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(0.5), Inches(0.5))
    tfIcon = txIcon.text_frame
    tfIcon.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tfIcon.paragraphs[0].add_run()
    run.text = str(slide_index)
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 标题文字
    txTitle = slide.shapes.add_textbox(Inches(1.3), Inches(0.3), Inches(10.5), Inches(0.8))
    tfTitle = txTitle.text_frame
    tfTitle.word_wrap = True
    _set_text(tfTitle, slide_data.get("title", ""), font_size=26, color=colors["primary"], bold=True)

    # 标题下划线
    _add_rect(slide, colors["accent"], Inches(1.3), Inches(1.15), Inches(1.5), Inches(0.05))

    # 页码
    txNum = slide.shapes.add_textbox(Inches(11.5), Inches(0.35), Inches(1.5), Inches(0.5))
    tfNum = txNum.text_frame
    _set_text(tfNum, f"{slide_index} / {total_slides}", font_size=11,
              color=colors["text_light"], alignment=PP_ALIGN.RIGHT)

    # 内容区域 - 卡片式布局
    content = slide_data.get("content", [])
    if content:
        # 内容卡片背景
        _add_rect(slide, colors["card_bg"], Inches(0.6), Inches(1.6), Inches(12.0), Inches(5.2), radius=0.05)

        # 内容要点
        for i, point in enumerate(content[:6]):  # 最多6个要点
            y = 1.85 + i * 0.8

            # 图标形状
            icon_shape = POINT_ICONS[i % len(POINT_ICONS)]
            _add_icon_shape(slide, icon_shape, colors["accent"], Inches(1.0), Inches(y + 0.08), Inches(0.3))

            # 要点文字
            txPoint = slide.shapes.add_textbox(Inches(1.5), Inches(y), Inches(10.5), Inches(0.7))
            tfPoint = txPoint.text_frame
            tfPoint.word_wrap = True
            _set_text(tfPoint, point, font_size=16, color=colors["text"])

    # 备注
    notes = slide_data.get("notes", "")
    if notes:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = notes

    # 底部装饰线
    _add_rect(slide, colors["accent"], Inches(0.6), Inches(7.0), Inches(12.0), Inches(0.02))

    return slide


def generate_end_slide_v2(prs, colors):
    """生成增强版结束页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 主背景
    _add_rect(slide, colors["primary"], Inches(0), Inches(0), Inches(13.333), Inches(7.5))

    # 装饰元素
    _add_circle(slide, colors["accent"], Inches(-0.5), Inches(5.5), Inches(2.5))
    _add_circle(slide, colors["accent"], Inches(11.0), Inches(-0.5), Inches(2.0))
    circle2 = _add_circle(slide, colors["accent2"], Inches(10.0), Inches(0.3), Inches(1.2))
    circle2.fill.fore_color.brightness = 0.2

    # 顶部装饰条
    _add_rect(slide, colors["accent"], Inches(0), Inches(0), Inches(13.333), Inches(0.08))

    # 中间装饰线
    _add_rect(slide, colors["accent"], Inches(5.0), Inches(2.8), Inches(3.333), Inches(0.04))

    # 谢谢
    txBox = slide.shapes.add_textbox(Inches(2), Inches(3.0), Inches(9.333), Inches(1.5))
    tf = txBox.text_frame
    _set_text(tf, "谢谢", font_size=48, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, alignment=PP_ALIGN.CENTER)

    # THANK YOU
    txBox2 = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(9.333), Inches(0.8))
    tf2 = txBox2.text_frame
    _set_text(tf2, "THANK YOU", font_size=18, color=RGBColor(0xCB, 0xD5, 0xE1), alignment=PP_ALIGN.CENTER)

    # 底部信息
    _add_rect(slide, colors["primary"], Inches(0), Inches(6.5), Inches(13.333), Inches(1.0))
    txInfo = slide.shapes.add_textbox(Inches(2), Inches(6.65), Inches(9.333), Inches(0.6))
    tfInfo = txInfo.text_frame
    _set_text(tfInfo, "Powered by MiMo AI  ·  智能演示文稿生成器",
              font_size=11, color=RGBColor(0x94, 0xA3, 0xB8), alignment=PP_ALIGN.CENTER)

    return slide


def generate_ppt(slides_data, template_id="blueprint", output_dir=None):
    """
    生成PPT文件

    Args:
        slides_data: dict, 包含 title 和 slides 列表
        template_id: str, 模板ID
        output_dir: str, 输出目录

    Returns:
        str: 生成的PPT文件路径
    """
    if output_dir is None:
        from config import OUTPUT_DIR
        output_dir = OUTPUT_DIR

    colors = COLOR_SCHEMES.get(template_id, COLOR_SCHEMES["blueprint"])
    title = slides_data.get("title", "演示文稿")
    slides = slides_data.get("slides", [])

    # Create presentation (16:9)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. 标题页
    subtitle = f"共 {len(slides)} 页 · AI 自动生成"
    generate_title_slide(prs, title, subtitle, colors)

    # 2. 目录页（2页以上才生成）
    if len(slides) >= 2:
        generate_toc_slide(prs, slides, colors)

    # 3. 内容页
    for i, slide_data in enumerate(slides):
        generate_content_slide_v2(prs, slide_data, i + 1, len(slides), colors)

    # 4. 结束页
    generate_end_slide_v2(prs, colors)

    # Save
    filename = f"ppt_{uuid.uuid4().hex[:8]}.pptx"
    filepath = os.path.join(output_dir, filename)
    prs.save(filepath)

    return filepath
